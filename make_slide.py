#!/usr/bin/env python3
import sys
import json
import os
import tempfile
import time
from copy import deepcopy
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RgbColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import requests

def clone_slide(prs, slide_to_clone):
    # Slaydni togri nusxalash
    blank_layout = slide_to_clone.slide_layout
    new_slide = prs.slides.add_slide(blank_layout)

    for shape in slide_to_clone.shapes:
        el = shape.element
        new_el = deepcopy(el)
        new_slide.shapes._spTree.insert_element_before(new_el, 'p:extLst')

    return new_slide

def fill_slide(prs, slide_idx, data, lang='uz'):
    # Slaydni malumotlar bilan toldirish
    slide = prs.slides[slide_idx]

    for shape in slide.shapes:
        # FIX: shape.name None bolishi mumkin
        nm = (shape.name or '').lower()

        if shape.has_text_frame and ('title' in nm) and ('footer' not in nm):
            tf = shape.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.text = data.get('title', '')
            p.font.size = Pt(28)
            p.font.bold = True

        elif shape.has_text_frame and ('content' in nm or 'body' in nm):
            tf = shape.text_frame
            tf.clear()
            bullets = data.get('bullets', [])
            for i, bullet in enumerate(bullets):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = f"• {bullet}"
                p.font.size = Pt(16)
                p.level = 0

def add_image_to_slide(prs, slide_idx, img_url, lang='uz'):
    # Rasm qoshish
    try:
        # FIX: tempfile ishlatish (Railway da /tmp yoq)
        img_path = os.path.join(tempfile.gettempdir(), f"slide_img_{slide_idx}_{int(time.time())}.jpg")

        response = requests.get(img_url, timeout=30)
        if response.status_code == 200:
            with open(img_path, 'wb') as f:
                f.write(response.content)

            slide = prs.slides[slide_idx]
            slide.shapes.add_picture(img_path, Inches(6), Inches(1.5), width=Inches(3))

            if os.path.exists(img_path):
                os.remove(img_path)
    except Exception as e:
        print(f"Image error: {e}")

def generate_presentation(data, output_path):
    # Asosiy funksiya - prezentatsiya yaratish
    topic = data.get('topic', 'Untitled')
    slides_data = data.get('slides', [])
    lang = data.get('lang', 'uz')
    user = data.get('user', {})
    template_id = data.get('template', 'template1')

    # Template yoli - Railway uchun process.cwd()
    template_path = os.path.join(os.getcwd(), f"{template_id}.pptx")

    if os.path.exists(template_path):
        prs = Presentation(template_path)
    else:
        prs = Presentation()

    # Muqova slaydi
    if len(prs.slides) > 0:
        fill_slide(prs, 0, {
            'title': topic,
            'bullets': [f"{user.get('name', '')} {user.get('surname', '')}".strip()]
        }, lang)

    # Kontent slaydlari
    for i, slide_data in enumerate(slides_data):
        if i + 1 < len(prs.slides):
            fill_slide(prs, i + 1, slide_data, lang)
        else:
            # Yangi slayd qoshish
            if len(prs.slides) > 1:
                new_slide = clone_slide(prs, prs.slides[1])
                fill_slide(prs, len(prs.slides) - 1, slide_data, lang)

    prs.save(output_path)
    print(f"Presentation saved: {output_path}")

if __name__ == '__main__':
    if len(sys.argv) < 3:
        print("Usage: python make_slide.py <input.json> <output.pptx>")
        sys.exit(1)

    json_path = sys.argv[1]
    output_path = sys.argv[2]

    with open(json_path, 'r', encoding='utf-8') as f:
        data = json.load(f)

    generate_presentation(data, output_path)
